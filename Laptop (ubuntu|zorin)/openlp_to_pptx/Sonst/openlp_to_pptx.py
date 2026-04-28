import os
import zipfile
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def extract_osz(osz_path, extract_to="/tmp/openlp_extract"):
    os.makedirs(extract_to, exist_ok=True)
    with zipfile.ZipFile(osz_path, 'r') as zip_ref:
        zip_ref.extractall(extract_to)
    return extract_to

def load_service_data(folder):
    for name in os.listdir(folder):
        if name.endswith(".osj") or name.endswith(".json"):
            with open(os.path.join(folder, name), "r", encoding="utf-8") as f:
                return json.load(f)
    raise FileNotFoundError("Keine service_data.osj/.json gefunden")

def clean_lyrics(text):
    text = re.sub(r'[\[\(\{].*?[\]\)\}]', '', text)  # Akkorde entfernen
    text = re.sub(r'^\s*[-•]\s*', '', text, flags=re.MULTILINE)  # Spiegelstriche entfernen
    text = text.replace('<br/>', '\n').replace('<br>', '\n')
    return text.strip()

def load_theme(theme_name):
    """
    Lädt ein OpenLP-Theme und gibt alle relevanten Werte für die PowerPoint-Erstellung zurück.
    """
    base = os.path.expanduser("~/.local/share/openlp/themes")
    theme_dir = os.path.join(base, theme_name)
    theme_file = os.path.join(theme_dir, f"{theme_name}.json")
    if not os.path.exists(theme_file):
        raise FileNotFoundError(f"Theme-Datei {theme_file} nicht gefunden")

    with open(theme_file, "r", encoding="utf-8") as f:
        raw_theme = json.load(f)

    # Hintergrundbildpfad
    bg_path = None
    if raw_theme.get("background_type") == "image":
        bg_file = raw_theme.get("background_filename", {}).get("parts", ["", ""])[-1]
        bg_path = os.path.join(theme_dir, bg_file)
        if not os.path.exists(bg_path):
            bg_path = None

    # Alle relevanten Keys ins Theme-Dict übernehmen
    theme = {
        "font_name": raw_theme.get("font_main_name", "Arial"),
        "font_size": raw_theme.get("font_main_size", 60),  # OpenLP-Größe
        "font_color": raw_theme.get("font_main_color", "#FFFFFF"),
        "bg_path": bg_path,
        "display_horizontal_align": raw_theme.get("display_horizontal_align", 2),
        "display_vertical_align": raw_theme.get("display_vertical_align", 2),
        "font_main_x": raw_theme.get("font_main_x", 10),
        "font_main_y": raw_theme.get("font_main_y", 10),
        "background_color": raw_theme.get("background_color", "#FFFFFF"),
    }

    # Debug-Ausgabe
    print("\n📂 GELADENES THEME:")
    for k, v in theme.items():
        print(f"  {k}: {v}")

    return theme

def hex_to_rgb(hexcolor):
    hexcolor = hexcolor.lstrip("#")
    return RGBColor(int(hexcolor[0:2], 16), int(hexcolor[2:4], 16), int(hexcolor[4:6], 16))

def add_footer_info(slide, song_title, verse_name, author_text, theme):
    """
    Fügt unten links auf der Folie Liedtitel, aktuelle Strophe und Autor ein.
    """
    if not author_text:
        author_text = "Unknown"
    if not verse_name:
        verse_name = ""

    text = f"{song_title} – {verse_name}\nDichter: {author_text}"

    left = Inches(0.5)
    top = Inches(theme.get("font_main_y", 10) / 72 + 6.5)  # ungefähr unten
    width = Inches(12)
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.clear()  # sicherstellen, dass nichts drinsteht
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(theme["font_size"] * 0.35)  # kleiner als Haupttext
    p.font.name = theme["font_name"]
    try:
        p.font.color.rgb = RGBColor.from_string(theme["font_color"].replace("#",""))
    except:
        pass
    p.alignment = PP_ALIGN.LEFT

def create_pptx(songs, theme, output_path):
    from pprint import pprint

    #print("\n📂 GELADENES THEME:")
    #pprint(theme)
    #print("\n---")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    font_name = theme.get("font_name", "Arial")
    font_color = hex_to_rgb(theme.get("font_color", "#FFFFFF"))
    scale_factor = 0.5
    font_size = Pt(theme.get("font_size", 48) * scale_factor)

    horiz_align = theme.get("display_horizontal_align", 2)
    vert_align = theme.get("display_vertical_align", 2)
    x_start = Inches(theme.get("font_main_x", 0) / 50)
    y_start = Inches(theme.get("font_main_y", 0) / 50)

    # Zuordnungen
    horiz_map = {
        0: PP_ALIGN.LEFT,
        1: PP_ALIGN.RIGHT,
        2: PP_ALIGN.CENTER
    }
    vert_map = {
        0: MSO_ANCHOR.TOP,
        2: MSO_ANCHOR.BOTTOM,
        1: MSO_ANCHOR.MIDDLE
    }

    for song in songs:
        # === Titelfolie ===
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if theme.get("bg_path") and os.path.exists(theme["bg_path"]):
            slide.shapes.add_picture(theme["bg_path"], 0, 0,
                                     width=prs.slide_width,
                                     height=prs.slide_height)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3),
                                         prs.slide_width - Inches(2),
                                         Inches(2))
        tf = txBox.text_frame
        tf.text = song["title"]
        p = tf.paragraphs[0]
        p.font.name = font_name
        p.font.size = Pt(theme.get("font_size", 48) * 0.7)
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # === Songfolien ===
        for verse_name, slide_text in song["slides_with_names"]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            if theme.get("bg_path") and os.path.exists(theme["bg_path"]):
                slide.shapes.add_picture(theme["bg_path"], 0, 0,
                                         width=prs.slide_width,
                                         height=prs.slide_height)

            x_start = Inches(theme.get("font_main_x", 10) / 72)
            y_start = Inches(theme.get("font_main_y", 10) / 72)

            txBox = slide.shapes.add_textbox(
                x_start,
                y_start,
                prs.slide_width - x_start * 2,
                prs.slide_height - y_start * 2
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = slide_text

            horiz_align = theme.get("display_horizontal_align", 1)
            vert_align = theme.get("display_vertical_align", 1)

            for p in tf.paragraphs:
                p.font.name = font_name
                p.font.size = font_size
                p.font.color.rgb = font_color
                p.alignment = horiz_map.get(horiz_align, PP_ALIGN.CENTER)
            tf.vertical_anchor = vert_map.get(vert_align, MSO_ANCHOR.MIDDLE)

            # === Footer mit Autor, Songtitel und Strophe ===
            add_footer_info(slide, song["title"], verse_name, song["author"], theme)

    prs.save(output_path)
    print(f"✅ Präsentation erstellt: {output_path}")

def main():
    import argparse
    parser = argparse.ArgumentParser(description="OpenLP OSZ/.OSJ → PowerPoint 16:9")
    parser.add_argument("osz_path", help="Pfad zur .osz-Datei")
    parser.add_argument("-o", "--output", default="output.pptx", help="Ausgabedatei (.pptx)")
    args = parser.parse_args()

    extract_folder = extract_osz(args.osz_path)
    data = load_service_data(extract_folder)
    theme_name = data[0].get("openlp_core", {}).get("service-theme", "Geo Purple")
    theme = load_theme(theme_name)

    songs = []
    for item in data:
        if "serviceitem" in item:
            s = item["serviceitem"]
            title = s["header"]["title"]
            # Autor: aus header/footer auslesen, fallback auf data[0]["authors"]
            footer_list = s["header"].get("footer", [])
            author = "Unknown"
            if len(footer_list) > 1:
                author = footer_list[1].replace("Written by: ", "")
            elif s["data"] and "authors" in s["data"][0]:
                author = s["data"][0]["authors"]

            slides_with_names = []
            for verse in s["data"]:
                raw_slide = verse.get("raw_slide", "").strip()
                verse_name = verse.get("verseTag", "") or verse.get("title", "")
                if raw_slide:
                    slides_with_names.append((verse_name, clean_lyrics(raw_slide)))
            songs.append({"title": title, "slides_with_names": slides_with_names, "author": author})

    create_pptx(songs, theme, args.output)

if __name__ == "__main__":
    main()

