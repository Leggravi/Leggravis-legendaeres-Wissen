#!/usr/bin/env python3
"""
openlp_to_pptx.py — Ultimates OpenLP OSZ → PowerPoint Konvertiertool
Schatten/Outline via Text-Duplikation (kein XML-Effect), clean-top-Layer
"""

import os, sys, math, zipfile, json, re, argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ──────────────────────────────────────────────────────
# LAYOUT
# ──────────────────────────────────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.5
TEXT_L  = 0.6
TEXT_R  = 0.6
TEXT_T  = 0.5
TEXT_MAX_H = SLIDE_H * 0.77 - TEXT_T   # ~5.28"

# ──────────────────────────────────────────────────────
# FARBE
# ──────────────────────────────────────────────────────
def hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def is_light(h):
    h = h.lstrip("#")
    r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
    return (0.299*r + 0.587*g + 0.114*b)/255 > 0.5

def contrast(h):
    return "#000000" if is_light(h) else "#ffffff"

# ──────────────────────────────────────────────────────
# SMARTER ZEILENUMBRUCH
# ──────────────────────────────────────────────────────
PUNCT_BREAK = set('.!?;:,–—-')
CONJ_DE     = {'und','oder','aber','doch','denn','weil','dass','wenn','als','wie',
               'so','auch','noch','schon','nur','gar','doch'}
CONJ_EN     = {'and','or','but','yet','for','nor','so','because','that','when',
               'as','if','though'}
CONJUNCTIONS = CONJ_DE | CONJ_EN

def smart_wrap(text, cpw):
    """Bricht text intelligent bei max cpw Zeichen um.
    Bevorzugt: nach Satzzeichen, vor Konjunktionen, sonst Wortgrenze."""
    if len(text) <= cpw:
        return [text]
    lines = []
    remaining = text
    while len(remaining) > cpw:
        seg = remaining[:cpw]
        best = -1

        # 1. Nach Satzzeichen (rückwärts in den letzten 40 % des Segments)
        for i in range(len(seg)-1, max(int(len(seg)*0.5), 0)-1, -1):
            if seg[i] in PUNCT_BREAK:
                best = i + 1
                break

        # 2. Vor Konjunktion (rückwärts)
        if best == -1:
            words = seg.split(' ')
            for j in range(len(words)-1, 0, -1):
                if words[j].lower().rstrip('.,;:!?') in CONJUNCTIONS:
                    best = len(' '.join(words[:j]))
                    break

        # 3. Letzte Wortgrenze
        if best == -1:
            best = seg.rfind(' ')
            if best == -1:
                best = cpw  # harter Umbruch

        lines.append(remaining[:best].rstrip())
        remaining = remaining[best:].lstrip()

    if remaining:
        lines.append(remaining)
    return lines

def cpw_for(pt, box_w_in):
    """Zeichen pro Zeile bei gegebener Schriftgröße."""
    return max(8, int(box_w_in / (0.55 * pt / 72.0)))

# ──────────────────────────────────────────────────────
# SCHRIFTGRÖSSE ANPASSEN
# ──────────────────────────────────────────────────────
def fit_pt(text, base_pt, box_w_in, box_h_in):
    LINE_SP = 1.4
    avail   = box_h_in * 72.0

    def total_vis(pt):
        cpw = cpw_for(pt, box_w_in)
        n = 0
        for ln in text.split('\n'):
            n += 1 if not ln.strip() else len(smart_wrap(ln, cpw))
        return n

    if total_vis(base_pt) * base_pt * LINE_SP <= avail:
        return base_pt
    lo, hi = 8.0, base_pt
    for _ in range(25):
        mid = (lo+hi)/2
        if total_vis(mid) * mid * LINE_SP <= avail: lo = mid
        else: hi = mid
    return max(8.0, lo)

# ──────────────────────────────────────────────────────
# TEXT IN SEITEN AUFTEILEN
# ──────────────────────────────────────────────────────
def split_pages(text, pt, box_w_in, max_lines):
    """
    Teilt text in Seiten mit max max_lines Zeilen.
    Trennt nie eine logische Zeile zwischen zwei Seiten (außer wenn nötig).
    """
    cpw   = cpw_for(pt, box_w_in)
    logical = text.split('\n')

    # Jede logische Zeile → Liste von Wrap-Sub-Zeilen
    groups = []
    for ll in logical:
        if not ll.strip():
            groups.append(('empty', ['']))
        else:
            groups.append(('text', smart_wrap(ll, cpw)))

    pages = []
    cur_lines = []
    cur_count = 0

    for kind, subs in groups:
        n = len(subs)
        if kind == 'empty':
            # Leerzeile: nur hinzufügen wenn noch Platz
            if cur_count < max_lines:
                cur_lines.append(('empty', subs))
                cur_count += 1
            continue

        if n > max_lines:
            # Logische Zeile länger als max_lines → muss aufgeteilt werden
            if cur_lines:
                pages.append(cur_lines)
                cur_lines, cur_count = [], 0
            for i in range(0, n, max_lines):
                pages.append([('chunk', subs[i:i+max_lines])])
        elif cur_count + n > max_lines:
            # Passt nicht mehr → neue Seite
            if cur_lines:
                pages.append(cur_lines)
            cur_lines = [('text', subs)]
            cur_count = n
        else:
            cur_lines.append(('text', subs))
            cur_count += n

    if cur_lines:
        pages.append(cur_lines)

    # Seiten als (text_str, continuation_flags) zurückgeben
    # continuation_flags[i] = True wenn Zeile i eine Fortsetzungszeile ist
    result = []
    for page in pages:
        lines = []
        is_cont = []
        for kind, subs in page:
            for j, s in enumerate(subs):
                lines.append(s)
                is_cont.append(j > 0)  # Fortsetzungszeile wenn j>0
        result.append((lines, is_cont))
    return result if result else [([text], [False])]

# ──────────────────────────────────────────────────────
# OUTLINE-BREITE SKALIERT NACH SCHRIFTGRÖSSE
# ──────────────────────────────────────────────────────
def outline_width(pt):
    """Skaliert: 1 pt bei kleiner Schrift, bis 6 pt bei großer."""
    return max(1.0, min(6.0, pt * 0.08))

def shadow_offset(pt, offset=0.04):
    """Versatz des Schatten-Duplikats in Inches."""
    return min(offset, pt * 0.0018)

# ──────────────────────────────────────────────────────
# EINZELNE TEXTBOX-SCHICHT
# ──────────────────────────────────────────────────────
NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
def atag(n): return f"{{{NS}}}{n}"

def _set_outline_xml(run, width_pt, color_hex):
    """Setzt <a:ln> direkt in <a:rPr>."""
    r_el = run._r
    rPr  = r_el.find(atag("rPr"))
    if rPr is None:
        rPr = etree.Element(atag("rPr"))
        r_el.insert(0, rPr)
    for x in rPr.findall(atag("ln")): rPr.remove(x)
    ln = etree.SubElement(rPr, atag("ln"))
    ln.set("w", str(int(width_pt * 12700)))
    sf  = etree.SubElement(ln, atag("solidFill"))
    clr = etree.SubElement(sf,  atag("srgbClr"))
    clr.set("val", color_hex.lstrip("#"))
    rPr.insert(0, ln)

H_MAP = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.RIGHT, 2: PP_ALIGN.CENTER}
V_MAP = {0: MSO_ANCHOR.TOP, 1: MSO_ANCHOR.MIDDLE, 2: MSO_ANCHOR.BOTTOM}

def _add_layer(slide, l, t, w, h,
               lines, is_cont,
               font_name, pt, font_hex,
               h_align, v_align,
               bold=False, outline_pt=0,
               tight_spacing=True):
    """
    Erstellt EINE Textbox-Schicht.
    is_cont[i]=True → Zeile i ist Fortsetzung → weniger Abstand nach oben.
    outline_pt > 0 → Umriss via XML.
    """
    tb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap       = False   # wir umbrechen manuell
    tf.auto_size       = None
    tf.vertical_anchor = v_align
    rgb = hex_rgb(font_hex)

    NORM_SP = Pt(pt * 0.20)   # Abstand zwischen logischen Zeilen
    CONT_SP = Pt(pt * 0.02)   # Abstand bei Fortsetzungszeile

    first = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = h_align
        # Zeilenabstand
        if tight_spacing and i > 0:
            p.space_before = CONT_SP if is_cont[i] else NORM_SP
        run = p.add_run()
        run.text           = line
        run.font.name      = font_name
        run.font.size      = Pt(pt)
        run.font.color.rgb = rgb
        run.font.bold      = bold
        if outline_pt > 0:
            _set_outline_xml(run, outline_pt, contrast(font_hex))

    return tb

# ──────────────────────────────────────────────────────
# HAUPT-TEXTBLOCK MIT ALLEN EFFEKT-EBENEN
# ──────────────────────────────────────────────────────
def add_text_block(slide, l, t, w, h,
                   lines, is_cont,
                   font_name, pt, font_hex,
                   h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.TOP,
                   bold=False, effects="both", shadow_offset_val=0.04):
    """
    effects: 'none' | 'shadow' | 'outline' | 'both'
    Schichten (von hinten nach vorne):
      shadow  → Duplikat leicht versetzt in Kontrastfarbe
      outline → Duplikat mit XML-Outline
      clean   → Originaler Text oben ohne Effekte
    """
    do_shadow  = effects in ("shadow", "both")
    do_outline = effects in ("outline", "both")
    ow = outline_width(pt)
    dx = dy = shadow_offset(pt, shadow_offset_val)

    if do_shadow:
        _add_layer(slide, l+dx, t+dy, w, h,
                   lines, is_cont, font_name, pt, contrast(font_hex),
                   h_align, v_align, bold=bold, outline_pt=0)

    if do_outline:
        _add_layer(slide, l, t, w, h,
                   lines, is_cont, font_name, pt, font_hex,
                   h_align, v_align, bold=bold, outline_pt=ow)

    # Saubare Top-Schicht (immer, außer kein Effekt ohne Ebene)
    _add_layer(slide, l, t, w, h,
               lines, is_cont, font_name, pt, font_hex,
               h_align, v_align, bold=bold, outline_pt=0)

# ──────────────────────────────────────────────────────
# HINTERGRUND
# ──────────────────────────────────────────────────────
def set_bg(slide, bg_path, bg_color="#000000"):
    if bg_path and os.path.exists(bg_path):
        pic = slide.shapes.add_picture(
            bg_path, 0, 0,
            width=Inches(SLIDE_W), height=Inches(SLIDE_H))
        sp = slide.shapes._spTree
        sp.remove(pic._element)
        sp.insert(2, pic._element)
    else:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_rgb(bg_color)

def black_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0,0,0)

# ──────────────────────────────────────────────────────
# ROLLEN-ÜBERSETZUNG
# ──────────────────────────────────────────────────────
ROLE_PAT = [
    (r'^words?\s*and\s*music\s*(?:by)?:?\s*', "Worte & Musik"),
    (r'^words?\s*&\s*music\s*(?:by)?:?\s*',   "Worte & Musik"),
    (r'^written\s*by:?\s*',                    "Worte & Musik"),
    (r'^words?\s*(?:by)?:?\s*',                "Worte"),
    (r'^music\s*(?:by)?:?\s*',                 "Musik"),
    (r'^translation\s*(?:by)?:?\s*',           "Übersetzung"),
    (r'^translated\s*(?:by)?:?\s*',            "Übersetzung"),
    (r'^arrangement\s*(?:by)?:?\s*',           "Arrangement"),
    (r'^copyright:?\s*',                        "©"),
]
def translate_credit(line):
    for pat, label in ROLE_PAT:
        m = re.match(pat, line, re.IGNORECASE)
        if m:
            rest = line[m.end():].strip()
            return f"{label}: {rest}" if rest else label
    return line

CCLI_RE = re.compile(
    r'^ccli\s*(?:song|licen[sc]e)?\s*(?:number|#|nr\.?)?\s*:?\s*(\d+)',
    re.IGNORECASE)

# ──────────────────────────────────────────────────────
# THEME LADEN
# ──────────────────────────────────────────────────────
def load_theme(name):
    base = os.path.expanduser("~/.local/share/openlp/themes")
    tf   = os.path.join(base, name, f"{name}.json")
    DEF  = dict(font_name="Arial", font_size=48, font_color="#FFFFFF",
                bg_path=None, bg_color="#000000", h_align=2, v_align=0)
    if not os.path.exists(tf):
        print(f"⚠️  Theme '{name}' nicht gefunden", file=sys.stderr)
        return DEF
    with open(tf, encoding="utf-8") as f:
        r = json.load(f)
    bg_path = None
    if r.get("background_type") == "image":
        parts = r.get("background_filename", {})
        fn    = parts.get("parts",["",""])[-1] if isinstance(parts,dict) else str(parts)
        cand  = os.path.join(base, name, fn)
        if os.path.exists(cand): bg_path = cand
    t = dict(
        font_name  = r.get("font_main_name",           "Arial"),
        font_size  = r.get("font_main_size",           48),
        font_color = r.get("font_main_color",          "#FFFFFF"),
        bg_path    = bg_path,
        bg_color   = r.get("background_color",         "#000000"),
        h_align    = r.get("display_horizontal_align", 2),
        v_align    = r.get("display_vertical_align",   0),
    )
    print(f"📂 Theme: {name}")
    for k,v in t.items():
        if k != "bg_path": print(f"   {k}: {v}")
    print(f"   bg_path: {t['bg_path'] or 'Farbe: '+t['bg_color']}")
    return t

# ──────────────────────────────────────────────────────
# OSZ ENTPACKEN
# ──────────────────────────────────────────────────────
def extract_osz(path):
    out = Path("/tmp/openlp") / Path(path).stem
    out.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(path) as z: z.extractall(out)
    return str(out)

def load_service(folder):
    for n in os.listdir(folder):
        if n.endswith(".osj") or n.endswith(".json"):
            with open(os.path.join(folder, n), encoding="utf-8") as f:
                return json.load(f)
    raise FileNotFoundError("Keine .osj/.json gefunden")

def clean(text):
    text = re.sub(r'[\[\(\{][^\]\)\}]*[\]\)\}]', '', text)
    text = re.sub(r'^\s*[-•]\s*', '', text, flags=re.MULTILINE)
    text = text.replace('<br/>','\n').replace('<br>','\n')
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()

def parse_songs(data):
    songs = []
    for item in data:
        if "serviceitem" not in item: continue
        s      = item["serviceitem"]
        header = s.get("header", {})
        title  = header.get("title", "Unbekannt")
        credits, ccli = [], None
        for line in header.get("footer", []):
            line = line.strip()
            if not line or line == title: continue
            m = CCLI_RE.match(line)
            if m: ccli = m.group(1); continue
            line = line.replace(" and ", " und ").replace(" And ", " Und ")
            credits.append(translate_credit(line))
        slides = []
        for v in s.get("data", []):
            raw = v.get("raw_slide","").strip()
            if raw:
                tag = v.get("verseTag","") or v.get("title","")
                slides.append((tag, clean(raw)))
        songs.append(dict(title=title, credits=credits, ccli=ccli, slides=slides))
    return songs

# ──────────────────────────────────────────────────────
# FOOTER-POSITION
# ──────────────────────────────────────────────────────
FOOTER_H = 1.25

def footer_pos(corner):
    """
    corner: 'bl' | 'br' | 'tl' | 'tr'
    Gibt (left, top, h_align, v_anchor) zurück.
    """
    is_bottom = corner in ('bl', 'br')
    is_right  = corner in ('br', 'tr')
    top  = (SLIDE_H - FOOTER_H - 0.1) if is_bottom else 0.1
    ha   = PP_ALIGN.RIGHT if is_right else PP_ALIGN.LEFT
    va   = MSO_ANCHOR.BOTTOM if is_bottom else MSO_ANCHOR.TOP
    return 0.35, top, SLIDE_W - 0.7, FOOTER_H, ha, va

# ──────────────────────────────────────────────────────
# FOOTER HINZUFÜGEN
# ──────────────────────────────────────────────────────
def add_footer(slide, song, verse_tag, th, show_ccli, ccli_ov,
               corner, effects, shadow_offset_val=0.04):
    ccli  = ccli_ov or song.get("ccli")
    lines = []
    tag_s = f" – {verse_tag}" if verse_tag else ""
    lines.append(f"{song['title']}{tag_s}")
    for c in song["credits"]: lines.append(c)
    if show_ccli and ccli: lines.append(f"CCLI: {ccli}")
    if not lines: return

    base_main = th["font_size"] * 0.5
    fpt = max(9.0, base_main * 0.28)

    l, t, w, h, ha, va = footer_pos(corner)
    is_cont = [False] * len(lines)
    add_text_block(slide, l, t, w, h,
                   lines, is_cont,
                   th["font_name"], fpt, th["font_color"],
                   h_align=ha, v_align=va,
                   effects=effects, shadow_offset_val=shadow_offset_val)

# ──────────────────────────────────────────────────────
# TEXT → LINES + IS_CONT (mit smartem Wrap)
# ──────────────────────────────────────────────────────
def text_to_lines(text, pt, box_w_in):
    """Gibt (lines, is_cont) zurück ohne Seitenaufteilung."""
    cpw = cpw_for(pt, box_w_in)
    lines, is_cont = [], []
    for ll in text.split('\n'):
        if not ll.strip():
            lines.append(''); is_cont.append(False)
        else:
            subs = smart_wrap(ll, cpw)
            for j, s in enumerate(subs):
                lines.append(s)
                is_cont.append(j > 0)
    return lines, is_cont

# ──────────────────────────────────────────────────────
# FOLIEN BAUEN
# ──────────────────────────────────────────────────────
def add_title_slide(prs, song, th, show_ccli, ccli_ov,
                    corner, effects, custom_pt=None, shadow_offset_val=0.04):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, th["bg_path"], th["bg_color"])
    pt = (custom_pt or (th["font_size"] * 0.5)) * 1.1
    lines, ic = text_to_lines(song["title"], pt, 11.33)
    add_text_block(slide, 1.0, 1.5, 11.33, 3.8,
                   lines, ic,
                   th["font_name"], pt, th["font_color"],
                   h_align=PP_ALIGN.CENTER, v_align=MSO_ANCHOR.MIDDLE,
                   bold=True, effects=effects, shadow_offset_val=shadow_offset_val)
    add_footer(slide, song, "", th, show_ccli, ccli_ov, corner, effects, shadow_offset_val)


def add_lyrics_slides(prs, song, vtag, text, th,
                      show_ccli, ccli_ov, corner, effects,
                      custom_pt, split_mode, max_lines, shadow_offset_val=0.04):
    """
    split_mode: 'shrink' | 'pages' | 'maxlines'
    """
    W = SLIDE_W - TEXT_L - TEXT_R
    base_pt = custom_pt or (th["font_size"] * 0.5)
    h_al = H_MAP.get(th.get("h_align",2), PP_ALIGN.CENTER)
    v_al = V_MAP.get(th.get("v_align",0), MSO_ANCHOR.TOP)
    text_top = TEXT_T + (FOOTER_H if corner in ('tl', 'tr') else 0)

    if split_mode == 'shrink':
        pt = fit_pt(text, base_pt, W, TEXT_MAX_H)
        lines, ic = text_to_lines(text, pt, W)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, th["bg_path"], th["bg_color"])
        add_text_block(slide, TEXT_L, text_top, W, TEXT_MAX_H,
                       lines, ic, th["font_name"], pt, th["font_color"],
                       h_align=h_al, v_align=v_al, effects=effects, shadow_offset_val=shadow_offset_val)
        add_footer(slide, song, vtag, th, show_ccli, ccli_ov, corner, effects, shadow_offset_val)

    elif split_mode in ('pages', 'maxlines'):
        ml = max_lines if split_mode == 'maxlines' else 9999
        pages = split_pages(text, base_pt, W, ml)
        total = len(pages)
        for i, (pg_lines, pg_cont) in enumerate(pages):
            tag = f"{vtag} ({i+1}/{total})" if total > 1 else vtag
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide, th["bg_path"], th["bg_color"])
            add_text_block(slide, TEXT_L, text_top, W, TEXT_MAX_H,
                           pg_lines, pg_cont,
                           th["font_name"], base_pt, th["font_color"],
                           h_align=h_al, v_align=v_al, effects=effects, shadow_offset_val=shadow_offset_val)
            add_footer(slide, song, tag, th, show_ccli, ccli_ov, corner, effects, shadow_offset_val)

# ──────────────────────────────────────────────────────
# PPTX ERZEUGEN
# ──────────────────────────────────────────────────────
def create_pptx(songs, th, out_path,
                black_slides, show_ccli, ccli_ov,
                corner, effects,
                custom_pt, split_mode, max_lines, shadow_offset_val=0.04):
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    for song in songs:
        if black_slides: black_slide(prs)
        add_title_slide(prs, song, th, show_ccli, ccli_ov,
                        corner, effects, custom_pt, shadow_offset_val)
        for vtag, txt in song["slides"]:
            add_lyrics_slides(prs, song, vtag, txt, th,
                              show_ccli, ccli_ov, corner, effects,
                              custom_pt, split_mode, max_lines, shadow_offset_val)
    prs.save(out_path)
    print(f"✅ Gespeichert: {out_path}")

# ──────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("osz_path")
    ap.add_argument("-o","--output",      default="output.pptx")
    ap.add_argument("--black-slides",     action="store_true")
    ap.add_argument("--show-ccli",        action="store_true")
    ap.add_argument("--ccli",             default=None)
    ap.add_argument("--corner",           default="bl",
                    choices=["bl","br","tl","tr"],
                    help="Footer-Ecke: bl/br/tl/tr")
    ap.add_argument("--effects",          default="both",
                    choices=["none","shadow","outline","both"])
    ap.add_argument("--font-size",        type=float, default=None,
                    help="Schriftgröße in pt (leer = aus Theme)")
    ap.add_argument("--split-mode",       default="shrink",
                    choices=["shrink","pages","maxlines"])
    ap.add_argument("--max-lines",        type=int, default=5,
                    help="Max Zeilen pro Folie bei --split-mode maxlines")
    ap.add_argument("--shadow-offset",    type=float, default=0.04,
                    help="Schattenversatz in Inches (default 0.04)")
    a = ap.parse_args()

    folder = extract_osz(a.osz_path)
    data   = load_service(folder)
    tname  = data[0].get("openlp_core",{}).get("service-theme","")
    th     = load_theme(tname)
    songs  = parse_songs(data)

    print(f"\n🎵 {len(songs)} Lied(er)")
    for s in songs:
        print(f"   • {s['title']} ({len(s['slides'])} Strophen)")

    create_pptx(songs, th, a.output,
                a.black_slides, a.show_ccli, a.ccli,
                a.corner, a.effects,
                a.font_size, a.split_mode, a.max_lines, a.shadow_offset)

if __name__ == "__main__":
    main()